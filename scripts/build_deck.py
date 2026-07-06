"""Generates the Bedrock Apex & Visualforce badge presentation deck.

Run:  python scripts/build_deck.py
Output: docs/Bedrock-AVF-Solution.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Salesforce-inspired palette
NAVY = RGBColor(0x03, 0x2D, 0x60)
BLUE = RGBColor(0x00, 0x70, 0xD2)
LIGHT = RGBColor(0xF3, 0xF6, 0xF9)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x54, 0x69, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x04, 0x84, 0x4B)  # green for "built/demo" tags

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def _set(tf_or_p, text, size, color, bold=False, italic=False, font="Segoe UI"):
    r = tf_or_p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = font
    return r


def add_rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def title_slide():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(4.55), SLIDE_W, Inches(0.06), BLUE)

    box = s.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _set(p, "Apex & Visualforce Domain Badge", 40, WHITE, bold=True)
    p2 = tf.add_paragraph()
    _set(p2, "Bedrock Programmatic Scenario - Proposed Solution", 22, RGBColor(0xBF, 0xD7, 0xF0))

    box2 = s.shapes.add_textbox(Inches(0.9), Inches(4.8), Inches(11.5), Inches(1.5))
    tf2 = box2.text_frame
    p = tf2.paragraphs[0]
    _set(p, "Divyendra Naidu", 20, WHITE, bold=True)
    p = tf2.add_paragraph()
    _set(p, "Senior Technical Architect  |  dnaidu@salesforce.com", 15, RGBColor(0xBF, 0xD7, 0xF0))
    notes(s, "Introduce yourself briefly. State that this session covers the Bedrock programmatic "
             "scenario: I will summarize the case, then walk primary requirements (with live demos "
             "in the org and code in the IDE), cover the design-only items, and finish with Q&A.")
    return s


def divider_slide(title, subtitle=""):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, Inches(0.9), Inches(3.0), Inches(1.4), Inches(0.09), BLUE)
    box = s.shapes.add_textbox(Inches(0.9), Inches(3.2), Inches(11.5), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    _set(tf.paragraphs[0], title, 36, WHITE, bold=True)
    if subtitle:
        p = tf.add_paragraph()
        _set(p, subtitle, 18, RGBColor(0xBF, 0xD7, 0xF0))
    return s


def content_slide(title, tag, blocks, note_text=""):
    """blocks: list of dicts: {'h': heading} or {'b': text, 'lvl': 0/1}"""
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.05), NAVY)
    # Title
    tbox = s.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(10.6), Inches(0.7))
    tf = tbox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(tf.paragraphs[0], title, 26, WHITE, bold=True)
    # Tag (e.g., LIVE DEMO / DESIGN / PSEUDO-CODE)
    if tag:
        tagcolor = ACCENT if tag in ("LIVE DEMO",) else BLUE
        tbx = s.shapes.add_textbox(Inches(11.0), Inches(0.32), Inches(2.0), Inches(0.42))
        ttf = tbx.text_frame
        ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ttf.paragraphs[0]
        pp.alignment = PP_ALIGN.RIGHT
        _set(pp, tag, 12, RGBColor(0x9F, 0xE8, 0xC8) if tag == "LIVE DEMO" else RGBColor(0xBF, 0xD7, 0xF0), bold=True)

    body = s.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.9))
    btf = body.text_frame
    btf.word_wrap = True
    first = True
    for blk in blocks:
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        if "h" in blk:
            p.space_before = Pt(8)
            p.space_after = Pt(3)
            _set(p, blk["h"], 16, BLUE, bold=True)
        else:
            lvl = blk.get("lvl", 0)
            p.level = lvl
            bullet = "•   " if lvl == 0 else "–   "
            _set(p, bullet + blk["b"], 13.5 if lvl == 0 else 12.5, DARK if lvl == 0 else GREY)
    if note_text:
        notes(s, note_text)
    return s


def code_slide(title, tag, intro, code_lines, note_text=""):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.05), NAVY)
    tbox = s.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(10.6), Inches(0.7))
    tf = tbox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(tf.paragraphs[0], title, 26, WHITE, bold=True)
    if tag:
        tbx = s.shapes.add_textbox(Inches(11.0), Inches(0.32), Inches(2.0), Inches(0.42))
        ttf = tbx.text_frame
        pp = ttf.paragraphs[0]
        pp.alignment = PP_ALIGN.RIGHT
        _set(pp, tag, 12, RGBColor(0xBF, 0xD7, 0xF0), bold=True)
    if intro:
        ibox = s.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5))
        _set(ibox.text_frame.paragraphs[0], intro, 13, GREY, italic=True)
    panel = add_rect(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(5.4), RGBColor(0x0B, 0x20, 0x3E))
    ctf = panel.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.2)
    ctf.margin_top = Inches(0.12)
    first = True
    for line in code_lines:
        p = ctf.paragraphs[0] if first else ctf.add_paragraph()
        first = False
        _set(p, line if line else " ", 11, RGBColor(0xCF, 0xE6, 0xFF), font="Consolas")
    if note_text:
        notes(s, note_text)
    return s


# ----------------------------------------------------------------------------
# Build deck
# ----------------------------------------------------------------------------
title_slide()

content_slide("Agenda", "", [
    {"b": "Introductions"},
    {"b": "Solution context and the case in my own words"},
    {"b": "Primary requirements (with live org demos and IDE code):"},
    {"b": "Client Plans (PDF + interactive map)", "lvl": 1},
    {"b": "Quote Processing (high-volume batch)", "lvl": 1},
    {"b": "Indirect Accounts component", "lvl": 1},
    {"b": "Sales-app Home Page map", "lvl": 1},
    {"b": "Dynamic Process Launcher component", "lvl": 1},
    {"b": "AI Client-Insight assistant (RAG)", "lvl": 1},
    {"b": "Data centricity (Data Cloud) and AI strategy"},
    {"b": "Secondary requirements"},
    {"b": "Q&A"},
], "Set expectations: ~35 minutes. Flag that four components are built and deployed and I will "
   "demo them live; the remaining items are design/pseudo-code as the scenario specifies.")

divider_slide("Introduction")

content_slide("Who Am I?", "", [
    {"h": "Background"},
    {"b": "Senior Technical Architect with deep Salesforce platform and programmatic experience."},
    {"b": "Hands-on across Apex, LWC, integration, and large-scale data processing."},
    {"h": "Why it matters for Bedrock"},
    {"b": "I balance declarative-first thinking with the right programmatic patterns."},
    {"b": "Focus on scalable, secure, governor-aware solutions and modern dev practices (SFDX, source-driven)."},
], "Personalize this slide with your own bio. Tie your experience to Bedrock's needs: speed for "
   "advisors, scale for data, and modern Lightning development.")

divider_slide("Solution Context")

content_slide("The Case in My Own Words", "", [
    {"h": "Who and what"},
    {"b": "Bedrock advisors need speed-to-answer on client calls: real-time quotes, account summaries, account plans (PDFs)."},
    {"b": "Center of Excellence mandates modern Salesforce dev methods and best practices."},
    {"h": "Core needs I extracted"},
    {"b": "Rich, printable Client Plans that now include an interactive map."},
    {"b": "Nightly, high-volume quote ingestion (up to 1M staged rows) processed natively."},
    {"b": "Reusable, configurable, secure Lightning components admins can place via App Builder."},
    {"b": "A 360 customer profile (Data Cloud) and an AI assistant grounded in that data (RAG)."},
    {"h": "Guiding principles"},
    {"b": "Use declarative where it fits; use Apex where control, scale, or logic demands it. Design for governor limits and security from the start."},
], "Show I understand the business problem before jumping to tech. Emphasize advisor speed, "
   "regulated-industry security, and scale as the recurring themes.")

divider_slide("Primary Requirements")

# Section 2 - Client Plans (design)
content_slide("Client Plans: PDF + Interactive Map", "DESIGN", [
    {"h": "Problem"},
    {"b": "Today a text-only Visualforce page is rendered to PDF and saved to the Account as a File for quarterly meetings."},
    {"b": "The CIO wants an interactive JavaScript map (countries of business + growth markets) on the Lightning record page AND inside the generated PDF - via a single button click, no print-screen."},
    {"h": "Key challenge"},
    {"b": "Visualforce renderAs='pdf' uses a server-side rendering engine that does NOT execute JavaScript, so an interactive JS map cannot draw itself in the PDF."},
    {"h": "Options"},
    {"b": "A) Render a static map image server-side (static map image service / pre-rendered SVG) and embed in the VF PDF. Pro: keeps native VF-to-PDF, simple, no new infra. Con: image is static.", "lvl": 1},
    {"b": "B) Headless-browser / Node rendering service (e.g., on Heroku) that runs the JS map and returns a PDF. Pro: pixel-perfect interactive map. Con: extra infra, callout, cost.", "lvl": 1},
    {"b": "C) Third-party doc generation (Conga/Drawloop, etc.). Pro: low-code. Con: licensing, less control.", "lvl": 1},
    {"h": "Recommendation"},
    {"b": "Option A for the PDF (capture the map state as a static image and embed it) + keep the live interactive map as an LWC on the record page. A single LWC button calls Apex to build the PDF (Blob) and saves it as a ContentVersion on the Account. Revisit Option B only if true interactivity in the PDF becomes mandatory."},
], "Lead with the core insight: the PDF engine has no JS runtime. Walk options and land on the "
   "hybrid: interactive LWC map on screen, static rendered image in the PDF, Apex generates and "
   "attaches the File. Mention governor-safe Blob handling and that code was not required here.")

# Section 3 - Quote Processing (pseudo code)
content_slide("Quote Processing: High-Volume Batch", "PSEUDO-CODE", [
    {"h": "Problem"},
    {"b": "A 'dumb' middleware loads up to 1M Healthcare_Quote_Staging__c rows (+ up to 10 detail rows each) at 11pm. Salesforce must create/update standard Quote + QuoteLineItem records linked to Accounts."},
    {"h": "Why Batch Apex"},
    {"b": "1M rows far exceeds synchronous limits. Batch Apex chunks work, gives each scope its own governor limits, and is restartable.", "lvl": 1},
    {"b": "Database.Batchable with a QueryLocator over staging; Database.Stateful to accumulate counts/errors; scope size tuned (e.g., 200) because of child detail processing and DML.", "lvl": 1},
    {"h": "Design elements"},
    {"b": "Initiation: Schedulable at 11pm chained AFTER the middleware load (or fired by a Platform Event when load completes)."},
    {"b": "Logic encapsulation: a service class does upsert-by-external-id so the same logic is unit-testable and reusable."},
    {"b": "Error handling: per-record try/catch + Database.upsert(allOrNone=false); failures written to an error log object, not lost; clean rows still commit."},
    {"b": "Idempotency: upsert on an External Id so re-runs update rather than duplicate."},
], "Justify async/batch explicitly (rubric asks for sync vs async vs batch). Stress bulk-safe DML, "
   "partial success, idempotent upsert by external id, and observability via an error log.")

code_slide("Quote Processing: Pseudo Code", "PSEUDO-CODE",
    "Native batch solution - process initiation, encapsulation, error handling, testing.",
    [
    "global class QuoteStagingBatch implements Database.Batchable<SObject>, Database.Stateful {",
    "    global Integer processed = 0, failed = 0;",
    "",
    "    global Database.QueryLocator start(Database.BatchableContext bc) {",
    "        return Database.getQueryLocator(",
    "            'SELECT Id, External_Id__c, Account_Ext_Id__c, ... FROM Healthcare_Quote_Staging__c');",
    "    }",
    "",
    "    global void execute(Database.BatchableContext bc, List<Healthcare_Quote_Staging__c> scope) {",
    "        // 1. Map staging -> standard Quote (match Account by external id)",
    "        // 2. Upsert Quotes by External_Id__c (allOrNone = false)",
    "        // 3. Query child Healthcare_QuoteDetail_Staging__c, build QuoteLineItems",
    "        // 4. Upsert line items; relate to parent Quote + PricebookEntry",
    "        // 5. for each SaveResult: if !success -> insert Quote_Load_Error__c, failed++",
    "        //    else processed++",
    "    }",
    "",
    "    global void finish(Database.BatchableContext bc) {",
    "        // notify ops (email/Platform Event); optionally chain detail batch or next run",
    "    }",
    "}",
    "// Initiation:  System.scheduleBatch(new QuoteStagingBatch(), 'Nightly Quotes', 0, 200);",
    "// Testing:     @isTest with Test.startTest/stopTest over a representative staging set,",
    "//              assert Quotes created/updated, bad rows logged, governor limits respected.",
    ],
    "Talk through scope size choice (200 due to child DML), Database.Stateful for tallies, "
    "allOrNone=false for partial success, external-id upsert for idempotency, and test approach.")

# Section 4A - Indirect Accounts (demo)
content_slide("Indirect Accounts Component", "LIVE DEMO", [
    {"h": "Problem"},
    {"b": "Contacts relate to multiple Accounts. Users want to see indirect Account relationships (via shared Contacts) directly on the Account record - admin-placeable and configurable."},
    {"h": "Solution"},
    {"b": "LWC 'indirectAccounts' on the Account record page, backed by Apex 'IndirectAccountController'."},
    {"b": "Apex walks AccountContactRelation: from this Account's Contacts out to the other Accounts those Contacts touch; aggregates the relationship-occurrence count.", "lvl": 1},
    {"b": "Returns Account name, industry, # contacts, lifetime value, relationship count in a sortable lightning-datatable.", "lvl": 1},
    {"h": "Config + best practices (the how & why)"},
    {"b": "App Builder design attributes: 'Number of items to show' (1-50) and 'Sort by' (relationship count | name) - so admins reconfigure with no code."},
    {"b": "@AuraEnabled(cacheable=true) for client-side caching and performance; WITH USER_MODE enforces FLS/sharing; bulk-safe single-pass aggregation."},
    {"h": "Demo"},
    {"b": "Open an Account -> show the list; change Sort by / item count in App Builder -> re-render."},
], "Live demo: open an Account with shared contacts; show sorting and the configurable properties "
   "in App Builder. In the IDE show the controller (AccountContactRelation traversal, USER_MODE, "
   "cacheable) and the targetConfig design attributes.")

# Section 1 - Home page map (demo)
content_slide("Sales-App Home Page Map", "LIVE DEMO", [
    {"h": "Problem"},
    {"b": "Add a Home-page component showing a pinned map of Bedrock HQ: 'The Landmark @ One Market, Suite 300, San Francisco, CA 94105', zoomed in but still showing the city."},
    {"h": "Solution"},
    {"b": "LWC 'bedrockHqMap' built in the IDE (SFDX) and deployed; uses the base lightning-map component."},
    {"b": "No Apex, no API key: lightning-map geocodes the address; marker + center set; zoom-level 16 balances street detail with city visibility.", "lvl": 1},
    {"h": "Options considered"},
    {"b": "lightning-map (base component) vs a custom Google Maps JS integration. Chose lightning-map: faster, secure, no key management, mobile-ready.", "lvl": 1},
    {"h": "Testing & security (considered, not implemented)"},
    {"b": "Would add Jest tests for the component and rely on LWS; no Apex means no sharing/FLS surface here."},
    {"h": "Demo"},
    {"b": "Show the component on the Sales app Home page; show the small, modern codebase in the IDE."},
], "Emphasize 'modern and efficient': a base component beats hand-rolled JS. Show js-meta target "
   "lightning__HomePage and how it's dropped on the Home page. Mention you considered Jest/LWS but "
   "did not implement, per the scenario.")

# Section 4B - Dynamic Process Launcher (demo)
content_slide("Dynamic Process Launcher", "LIVE DEMO", [
    {"h": "Problem"},
    {"b": "Regional Sales Managers define dynamic, attribute-driven processes. Admins drop a component on any Opportunity (and other objects - Support/HR/Marketing) that renders chevrons or buttons; each opens a modal for a specific process. Must be secure and look native."},
    {"h": "Solution"},
    {"b": "Generic LWC 'processLauncher' on the record page; reads @api recordId/objectApiName; renders an SLDS path (chevron) or button group from configuration."},
    {"b": "Configuration lives in Custom Metadata 'Process_Action__mdt' (object, label, order, style, icon, target) - Sales Managers add/reorder steps with NO code, and it works on any object.", "lvl": 1},
    {"b": "Apex 'ProcessLauncherController.getActions(objectApiName)' returns the active, sorted actions WITH USER_MODE.", "lvl": 1},
    {"h": "Security (regulated industry)"},
    {"b": "Lightning Web Security isolates the component; it does not broadcast record data on shared message channels, so unauthorized AppExchange components cannot read it. Public APIs are explicit (@api) and minimal."},
    {"h": "Look & feel + demo"},
    {"b": "Built entirely with SLDS so it matches Lightning Experience. Demo: chevrons on Opportunity, buttons on Case, all from the same component + metadata."},
], "This is the strongest 'platform craft' demo: one component, metadata-driven, reusable across "
   "objects, secure by design (LWS, no shared-channel leakage), SLDS-native. Show the CMDT records "
   "and how adding a row changes the UI with no deploy.")

# Section 7 - AI / RAG (demo)
content_slide("AI Client-Insight Assistant (RAG)", "LIVE DEMO", [
    {"h": "Problem"},
    {"b": "Advisors want instant, natural-language answers about a client during calls, grounded in real data - not hallucinations."},
    {"h": "Solution: RAG implemented in Apex + LWC"},
    {"b": "Retrieve: Apex 'ClientInsightService' queries Account, recent Opportunities, Cases, and attached files (account plans) WITH USER_MODE, bulk-safe with row limits.", "lvl": 1},
    {"b": "Augment: builds a grounded, size-bounded prompt (guards prompt/heap limits) and injects the running user's locale for language control.", "lvl": 1},
    {"b": "Generate: an 'LlmProvider' interface - live model via Named Credential, with automatic fallback to a deterministic provider so it always works.", "lvl": 1},
    {"h": "Engineering for trust"},
    {"b": "Exposed as @AuraEnabled (LWC) and @InvocableMethod (so the Section 7 Screen Flow gains Apex logging, error handling, and dynamic language it can't do alone)."},
    {"b": "Hallucination mitigation: answer ONLY from retrieved context; AI disclaimer in the UI; graceful error handling."},
    {"h": "Config + demo"},
    {"b": "App Builder attributes (card title, default question). Demo on an Account: ask a question / use a canned prompt; show grounded answer + the controller in the IDE."},
], "Frame the three RAG steps clearly (retrieve/augment/generate). Stress Apex value-add over "
   "Flow-only: locale-aware language, logging, error handling, input pre-processing. Note the "
   "provider interface = OOP/polymorphism and a reliable demo via the mock fallback.")

# Section 6 - Data Cloud (design)
content_slide("Data Centricity: Customer 360 (Data Cloud)", "DESIGN", [
    {"h": "Goal"},
    {"b": "Unify CRM, the quote back-end, service/support, and marketing/ERP into one 360 profile with identity resolution."},
    {"h": "Approach"},
    {"b": "Ingest via connectors / Ingestion API (not CSV) into Data Cloud DMOs; map sources to the Customer 360 data model consistently."},
    {"b": "Identity resolution rules merge duplicates; conflict handling via source precedence / last-updated-wins rules."},
    {"b": "Refresh near-real-time or via efficient nightly (11pm) ingestion of changed records since last run."},
    {"h": "Surfacing + scale"},
    {"b": "Show the unified profile in Sales Cloud via Data Cloud-related components or a custom LWC when advisors open an Account/Contact."},
    {"b": "Any write-back to core Salesforce must be bulk-safe; respect governor limits; robust service with retry."},
    {"b": "This unified profile becomes the richest retrieval source for the RAG assistant."},
], "Connect this to the AI slide: Data Cloud is the future retrieval layer for RAG. Emphasize "
   "connector-first, identity resolution, conflict rules, nightly+streaming, bulk-safe write-back.")

content_slide("Optional: High-Level Data Model", "", [
    {"h": "Core standard objects"},
    {"b": "Account 1--* Contact (and Account *--* Contact via AccountContactRelation for indirect relationships)."},
    {"b": "Account 1--* Opportunity 1--* OpportunityLineItem;  Opportunity --> Order (Closed Won sync)."},
    {"b": "Account 1--* Quote 1--* QuoteLineItem  (loaded from staging)."},
    {"h": "Custom / staging"},
    {"b": "Healthcare_Quote_Staging__c 1--* Healthcare_QuoteDetail_Staging__c (assumed; ingestion staging)."},
    {"b": "Process_Action__mdt (custom metadata) drives the Dynamic Process Launcher per object."},
    {"b": "Quote_Load_Error__c (assumed) captures batch failures for observability."},
    {"h": "Data Cloud"},
    {"b": "Source objects -> DMOs -> unified Individual/Account profile (identity resolution)."},
], "Optional slide. Use it to anchor assumptions: AccountContactRelation for indirect accounts, "
   "staging->Quote, CMDT for the launcher, and the Data Cloud unification.")

divider_slide("Secondary Requirements", "Code not required - design and recommendations")

content_slide("Secondary Requirements", "DESIGN", [
    {"h": "Orders sync on Closed Won (scale)"},
    {"b": "When Opportunity is Closed Won, create an Order, sync line items, set 'Ready for Integration', publish a Platform Event to the ESB. With thousands of line items, do the sync asynchronously (Queueable/Batch) and bulk-safe; event-based hand-off decouples from the Sales Order system."},
    {"h": "Extensibility - coordinating event-driven logic"},
    {"b": "Consolidate per-object logic: one trigger per object + a handler/Trigger Actions Framework (programmatic), with explicit, ordered execution. Declarative option: Flow + Flow Orchestration. Recommend a single framework so order is known and logic isn't scattered."},
    {"h": "Enable/disable customization (data loads)"},
    {"b": "Hierarchy Custom Setting and/or Custom Permissions to toggle automation by user/profile, plus a global 'bypass' switch so large data loads can suspend logic, then re-apply after load. Declarative + programmatic checks honor the same switch."},
    {"h": "Opportunity name pre-population (Date + Time, before Save)"},
    {"b": "Must populate before save on the standard page, from both the Account related list and the Opportunity tab. A before-save Flow runs after save, so override the New action with a lightweight LWC/quick action (or default-value logic) that pre-fills the name with current date+time."},
], "Each of these is design-only per the scenario. Hit the rubric points: async/batch + Platform "
   "Event for scale; a single trigger framework for ordered logic; custom settings/permissions with "
   "a bypass for data loads; and why before-save Flow can't satisfy the 'before Save' name pre-fill.")

divider_slide("Thank You", "Questions?")

out_dir = os.path.join(os.path.dirname(__file__), "..", "docs")
out_path = os.path.normpath(os.path.join(out_dir, "Bedrock-AVF-Solution.pptx"))
prs.save(out_path)
print("Saved:", out_path, "slides:", len(prs.slides._sldIdLst))
