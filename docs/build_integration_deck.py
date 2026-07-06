"""Build the Bedrock Integration Architecture badging deck (.pptx).

Produces an editable PowerPoint with:
  - Title / Agenda / Who Am I
  - Slide 1: Problem statement
  - Slide 2: Solution architecture (native shape diagram) + assumptions
  - Slide 3: Interface list (native table)
  - Slide 4: Integration security
  - Slide 5: Summary & trade-offs
  - Thank you
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- palette ----
SF_BLUE   = RGBColor(0x01, 0x76, 0xD3)   # Salesforce blue
MULE_BLUE = RGBColor(0x00, 0x6D, 0x9C)   # MuleSoft teal-blue
GREEN     = RGBColor(0x2E, 0x84, 0x4A)   # channels
NAVY      = RGBColor(0x03, 0x2D, 0x60)   # headers / backends
SLATE     = RGBColor(0x44, 0x4A, 0x59)   # backend boxes
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xF3, 0xF4, 0xF6)
GREYTX    = RGBColor(0x32, 0x32, 0x32)
LINE_GREY = RGBColor(0x9A, 0xA0, 0xAA)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, x, y, w, h, text, fill, font=12, bold=True, fg=WHITE,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.CENTER, line=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.color.rgb = fill
    else:
        sp.line.color.rgb = line
    sp.line.width = Pt(1)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(font)
        r.font.bold = bold if i == 0 else False
        r.font.color.rgb = fg
        r.font.name = "Segoe UI"
    return sp


def textbox(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, sz, bd, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bd
        r.font.color.rgb = col
        r.font.name = "Segoe UI"
    return tb


def title_bar(s, title, subtitle=None):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE
    r.font.name = "Segoe UI Semibold"
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0xBF, 0xD6, 0xF2)
        r2.font.name = "Segoe UI"


def connect(s, a, b, color=LINE_GREY, width=1.5, dash=False):
    """Straight connector between centers of shapes a and b (edge-aware)."""
    ax, ay = a.left + a.width // 2, a.top + a.height // 2
    bx, by = b.left + b.width // 2, b.top + b.height // 2
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax, ay, bx, by)
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    cn.shadow.inherit = False
    if dash:
        ln = cn.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(d)
    # send connectors behind boxes
    sp = cn._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(3, sp)
    return cn


def bullets(s, x, y, w, h, items, size=14, color=GREYTX, gap=6):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            txt, lvl, bold = it
        else:
            txt, lvl, bold = it, 0, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = ("\u2022 " if lvl == 0 else "\u2013 ") + txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Segoe UI"
    return tb


# ============================================================
# TITLE
# ============================================================
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.55), SW, Inches(0.08))
band.fill.solid(); band.fill.fore_color.rgb = SF_BLUE; band.line.fill.background(); band.shadow.inherit = False
textbox(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.4), [
    ("Bedrock Financial Services", 40, True, WHITE),
    ("Integration Architecture", 32, True, RGBColor(0x7F, 0xB8, 0xEE)),
    ("Master-data integration across CRM, customer & position masters,", 16, False, LIGHT),
    ("core banking, credit agencies and the enterprise data lake", 16, False, LIGHT),
])
textbox(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.8), [
    ("First Name Last Name  \u00b7  Integration Architect  \u00b7  email@salesforce.com", 13, False, RGBColor(0xBF, 0xD6, 0xF2)),
])

# ============================================================
# AGENDA
# ============================================================
s = slide(); title_bar(s, "Agenda")
bullets(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5), [
    "Introduction (elevator pitch)",
    "Problem statement",
    "Proposed solution \u2014 architecture & assumptions",
    "Interface list (patterns, APIs, security, error handling)",
    "Integration security",
    "Summary & trade-offs",
], size=18, gap=12)

# ============================================================
# WHO AM I
# ============================================================
s = slide(); title_bar(s, "Who Am I?")
bullets(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5), [
    "Integration / Salesforce architect.",
    "Experience designing API-led, event-driven integrations for regulated",
    "financial-services platforms at scale.",
    "Focus: choosing the right integration pattern, API, and security model per use case.",
], size=18, gap=10)

# ============================================================
# SLIDE 1 - PROBLEM STATEMENT
# ============================================================
s = slide(); title_bar(s, "High-Level Problem Statement", "Master-data fragmentation across many heterogeneous systems")
bullets(s, Inches(0.6), Inches(1.25), Inches(7.2), Inches(6.0), [
    ("Master-data problem across customers and securities/products.", 0, True),
    "Customer Master = Hadoop (EA-owned, partly outsourced; real-time + batch).",
    "Position Master = separate system; business wants real-time access in Salesforce.",
    "Salesforce = master for Chatter, Opportunities, Leads, Cases.",
    ("The challenge", 0, True),
    "Connect Marketing Cloud, Drupal portal/mobile, credit agencies, core banking,",
    "the masters, legacy RockStar and a data lake \u2014 each with different latency,",
    "volume and security profiles \u2014 using the correct integration pattern per case.",
    "Securely (mutual auth), centrally monitored, and as reusable assets.",
], size=13.5, gap=6)
# volumes panel
box(s, Inches(8.1), Inches(1.4), Inches(4.6), Inches(0.55), "Volumes / Scale", NAVY, font=15)
bullets(s, Inches(8.25), Inches(2.05), Inches(4.4), Inches(4.5), [
    "5M retail customers",
    "~2 accounts each (checking/savings + credit card)",
    "~10 transactions / account / day  (~100M/day)",
    "Multi-petabyte historical data warehouse",
    ">100,000 transaction-history calls/hour at peak",
    "Millions of loan-processing events / day",
    "150 marketing/sales reps + 500 service reps",
], size=13, gap=8, color=GREYTX)

# ============================================================
# SLIDE 2 - SOLUTION ARCHITECTURE (DIAGRAM) + ASSUMPTIONS
# ============================================================
s = slide(); title_bar(s, "Solution Architecture \u2014 API-led, Event-driven on MuleSoft")

# --- column headers ---
box(s, Inches(0.25), Inches(1.15), Inches(2.5), Inches(0.4), "CHANNELS", GREEN, font=11)
box(s, Inches(3.05), Inches(1.15), Inches(3.0), Inches(0.4), "SALESFORCE PLATFORM", SF_BLUE, font=11)
box(s, Inches(6.35), Inches(1.15), Inches(2.6), Inches(0.4), "INTEGRATION BACKBONE", MULE_BLUE, font=11)
box(s, Inches(9.25), Inches(1.15), Inches(3.85), Inches(0.4), "SYSTEMS OF RECORD / EXTERNAL", SLATE, font=11)

# --- channels (left) ---
ch_mc = box(s, Inches(0.25), Inches(1.75), Inches(2.5), Inches(0.75), "Marketing Cloud\n(campaigns)", GREEN, font=11)
ch_dr = box(s, Inches(0.25), Inches(2.75), Inches(2.5), Inches(0.9), "Drupal Portal\n(self-service + forms)", GREEN, font=11)
ch_mo = box(s, Inches(0.25), Inches(3.9), Inches(2.5), Inches(0.75), "Hybrid Mobile App", GREEN, font=11)
ch_tr = box(s, Inches(0.25), Inches(4.9), Inches(2.5), Inches(0.9), "External Trading\nSystem (UI mashup)", GREEN, font=11)

# --- Salesforce (center) ---
sf = box(s, Inches(3.05), Inches(1.75), Inches(3.0), Inches(4.05),
         "Salesforce", SF_BLUE, font=16)
sf.text_frame.vertical_anchor = MSO_ANCHOR.TOP
bullets(s, Inches(3.15), Inches(2.35), Inches(2.85), Inches(3.4), [
    "Sales & Service Cloud",
    "Headless Experience Cloud",
    "Platform Events / CDC",
    "Salesforce Connect (OData)",
    "Apex / Flow / LWC",
    "Composite & Pub-Sub APIs",
], size=11.5, gap=7, color=WHITE)

# --- MuleSoft backbone (center-right) ---
mule = box(s, Inches(6.35), Inches(1.75), Inches(2.6), Inches(4.05),
           "MuleSoft\nAnypoint", MULE_BLUE, font=15)
mule.text_frame.vertical_anchor = MSO_ANCHOR.TOP
bullets(s, Inches(6.45), Inches(2.55), Inches(2.45), Inches(3.2), [
    "System / Process /",
    "Experience APIs",
    "mTLS + OAuth/OIDC",
    "Central monitoring",
    "Reusable templates",
    "Error queue (DLQ)",
], size=11, gap=6, color=WHITE)

# --- systems of record (right) ---
sysx = Inches(9.25); sysw = Inches(3.85)
sy_cm = box(s, sysx, Inches(1.75), sysw, Inches(0.55), "Customer Master \u2014 Hadoop", SLATE, font=11)
sy_pm = box(s, sysx, Inches(2.42), sysw, Inches(0.55), "Position Master (real-time)", SLATE, font=11)
sy_cb = box(s, sysx, Inches(3.09), sysw, Inches(0.55), "Core Banking \u2014 OData", SLATE, font=11)
sy_ca = box(s, sysx, Inches(3.76), sysw, Inches(0.55), "Credit Agencies \u00d73+ (REST/OIDC)", SLATE, font=11)
sy_rs = box(s, sysx, Inches(4.43), sysw, Inches(0.55), "Legacy RockStar (migrate + sync)", SLATE, font=11)
sy_dl = box(s, sysx, Inches(5.10), sysw, Inches(0.55), "Enterprise Data Lake (CDC)", SLATE, font=11)

# --- connectors ---
connect(s, ch_mc, sf, color=GREEN, width=1.5)
connect(s, ch_dr, sf, color=GREEN, width=1.5)
connect(s, ch_mo, sf, color=GREEN, width=1.5)
connect(s, ch_tr, sf, color=GREEN, width=1.5, dash=True)
connect(s, sf, mule, color=SF_BLUE, width=2.5)
for b in (sy_cm, sy_pm, sy_cb, sy_ca, sy_rs, sy_dl):
    connect(s, mule, b, color=MULE_BLUE, width=1.5)
# Salesforce Connect virtualization (direct dashed to core banking)
connect(s, sf, sy_cb, color=RGBColor(0xC9, 0x6A, 0x00), width=1.25, dash=True)

# legend + assumptions strip
textbox(s, Inches(0.25), Inches(5.95), Inches(8.6), Inches(0.35), [
    ("Solid = via MuleSoft (monitored, mTLS, reusable)   \u00b7   Orange dashed = Salesforce Connect data virtualization   \u00b7   Blue = event/API bridge", 10, False, GREYTX),
])
box(s, Inches(0.25), Inches(6.35), Inches(2.0), Inches(0.4), "Key Assumptions", NAVY, font=12)
bullets(s, Inches(2.4), Inches(6.32), Inches(10.7), Inches(1.0), [
    "MuleSoft Anypoint licensed; certs/identities by EA  \u00b7  Customer Master = real-time+batch; Position Master callable  \u00b7  Core banking exposes OData",
    "Agencies = REST + OpenID Connect  \u00b7  Data lake can subscribe to a stream; PII list agreed  \u00b7  No Communities UI (headless Experience Cloud only)",
], size=10.5, gap=4, color=GREYTX)

# ============================================================
# SLIDE 3 - INTERFACE LIST (TABLE)
# ============================================================
s = slide(); title_bar(s, "Interface List", "Pattern \u00b7 Salesforce API \u00b7 Security \u00b7 Error handling per interface")

headers = ["ID", "Source \u2192 Target", "Content", "Pattern", "SF API / Mechanism", "Security", "Error handling"]
rows = [
    ["INT-01", "Marketing Cloud \u2192 SF", "Leads near-real-time to swarm", "Remote Call-In + UI Update", "Platform Event / Pub-Sub; empApi", "OAuth + mTLS", "Replay by replayId; DLQ"],
    ["INT-02", "Drupal \u2192 SF", "Capture prospect", "Remote Call-In (Req-Reply)", "REST API", "OAuth + mTLS", "Sync error; retry 5xx"],
    ["INT-03", "Rules engine \u2192 SF", "Daily call-list load", "Batch Data Sync", "Bulk API 2.0", "Named Cred + mTLS", "Failed-row reprocess"],
    ["INT-04", "SF UI \u2192 map", "Call list on a map", "UI integration", "lightning-map / Maps", "Key in Named Cred", "'Address not found'"],
    ["INT-05", "SF \u2192 Manager UI", "High-value call closed (live)", "UI Update on Data Change", "Platform Event + empApi", "Sharing/FLS", "Resubscribe + replay"],
    ["INT-06", "SF \u2192 Credit agencies", "Credit checks \u00d73+, \u22641 min, partial retry", "RPC Request & Reply", "Apex Continuation / Flow HTTP Callout \u2192 Mule fan-out", "OpenID Connect + mTLS", "'Agency XYZ down'; resubmit retries only failed"],
    ["INT-07", "SF \u2192 Hadoop + Position Master", "Registered customer near-real-time", "RPC Fire & Forget", "Platform Event \u2192 Mule fan-out", "mTLS", "Retry + DLQ \u2192 admin queue"],
    ["INT-08", "Master \u2192 Experience Cloud", "Create/activate user + invite", "Remote Call-In", "REST + createExternalUser", "OAuth + mTLS", "Retry queue"],
    ["INT-09", "Drupal \u2194 Experience Cloud", "Real-time self-service as end-user", "Remote Call-In (Req-Reply)", "Connect REST / headless APIs", "End-user OAuth \u2014 no API user", "401 \u2192 re-auth"],
    ["INT-10/12", "Drupal \u2192 SF", "Contact info / raise inquiry", "Remote Call-In", "REST API (end-user)", "End-user OAuth + mTLS", "Validation error to UI"],
    ["INT-11", "SF \u2192 Core banking", "Transaction history >100k/hr", "Data Virtualization", "Salesforce Connect (OData) + Mule cache", "Named Cred + mTLS", "Circuit breaker; cached"],
    ["INT-13", "SF \u2192 Customer Master", "Propagate updates", "RPC Fire & Forget", "Platform Event / CDC \u2192 Mule", "mTLS", "Retry + DLQ"],
    ["INT-14", "Mobile \u2192 SF", "Same functions as portal", "Remote Call-In", "REST / GraphQL", "End-user OAuth (PKCE) + mTLS", "Offline retry"],
    ["INT-15", "Drupal/Mule \u2192 SF", "Loan app all-or-nothing", "Remote Call-In (transactional)", "Composite API (allOrNone)", "OAuth + mTLS", "Whole-graph rollback"],
    ["INT-16", "SF \u2194 Trading system", "No swivel + auto-log", "UI mashup + RPC", "Canvas/iframe + SSO; log callout", "SSO + mTLS", "Log retry queue"],
    ["INT-17", "SF \u2192 ESB", "Millions of loan events/day", "RPC Fire & Forget (scale)", "High-Volume PE / Pub-Sub API", "OAuth + mTLS", "Replay + DLQ; monitor"],
    ["INT-18", "RockStar \u2194 SF", "Migration + 3-mo sync", "Batch + near-real-time sync", "Bulk API 2.0 + CDC \u2194 Mule", "mTLS", "Bulk error rows; conflict rules"],
    ["INT-19", "SF \u2192 Data Lake", "Non-PII; ins/upd/del; history, changed-fields-only", "RPC Fire & Forget (event-driven)", "Change Data Capture \u2192 Mule", "mTLS + PII filter", "Replay; gap detect; DLQ"],
]

n_rows = len(rows) + 1
n_cols = len(headers)
tbl_left, tbl_top = Inches(0.18), Inches(1.18)
tbl_w, tbl_h = Inches(13.0), Inches(6.1)
gtbl = s.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h).table

# disable banated styling glow by setting column widths
col_w = [0.72, 1.85, 2.05, 1.7, 2.35, 1.55, 1.78]  # inches, sums ~12.0... scale
scale = 13.0 / sum(col_w)
for i, w in enumerate(col_w):
    gtbl.columns[i].width = Inches(w * scale)

# header row
for c, htxt in enumerate(headers):
    cell = gtbl.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    cell.margin_left = Inches(0.04); cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = htxt
    r.font.size = Pt(8.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Segoe UI"

# body rows
for ri, row in enumerate(rows, start=1):
    band_col = LIGHT if ri % 2 == 0 else WHITE
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = band_col
        cell.margin_left = Inches(0.04); cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.005); cell.margin_bottom = Inches(0.005)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = val
        r.font.size = Pt(7.5)
        r.font.bold = (ci == 0)
        r.font.color.rgb = NAVY if ci == 0 else GREYTX
        r.font.name = "Segoe UI"

# row heights
gtbl.rows[0].height = Inches(0.3)
for ri in range(1, n_rows):
    gtbl.rows[ri].height = Inches(0.32)

# ============================================================
# SLIDE 4 - INTEGRATION SECURITY
# ============================================================
s = slide(); title_bar(s, "Integration Security", "Mutual auth on every interface \u00b7 end-user (not API user) for self-service")
sec = [
    ("Transport", "Mutual TLS (two-way SSL) on ALL interfaces \u2014 client + server certs; enforced centrally at MuleSoft API Manager and via Named Credentials."),
    ("Outbound (SF \u2192 external)", "Named + External Credentials; OpenID Connect / OAuth 2.0 to credit agencies; per-system credentials at Mule \u2014 no secrets in code."),
    ("Inbound (external \u2192 SF)", "OAuth 2.0 \u2014 client-credentials for server-to-server (Drupal form); Authorization Code / PKCE for end-user & mobile."),
    ("End-user, not API user", "Drupal & mobile call Experience Cloud with the END-USER's token \u2014 honoring per-user sharing & FLS. Meets the BFS security mandate (INT-09)."),
    ("Data protection", "Shield Platform Encryption for PII at rest; PII EXCLUDED from the data-lake CDC stream; field filtering in Mule."),
    ("Governance & monitoring", "Least-privilege connected apps / permission sets; central policy, threat protection & audit at API Manager + Salesforce Event Monitoring."),
]
yy = Inches(1.35)
for label, body in sec:
    box(s, Inches(0.5), yy, Inches(3.0), Inches(0.78), label, MULE_BLUE, font=12, align=PP_ALIGN.LEFT)
    textbox(s, Inches(3.7), yy + Inches(0.02), Inches(9.2), Inches(0.78),
            [(body, 12.5, False, GREYTX)], anchor=MSO_ANCHOR.MIDDLE)
    yy = yy + Inches(0.93)

# ============================================================
# SLIDE 5 - SUMMARY & TRADE-OFFS
# ============================================================
s = slide(); title_bar(s, "Summary & Trade-offs")
box(s, Inches(0.5), Inches(1.25), Inches(6.0), Inches(0.45), "Why MuleSoft \u2014 license justified", NAVY, font=14, align=PP_ALIGN.LEFT)
bullets(s, Inches(0.55), Inches(1.8), Inches(6.1), Inches(5.2), [
    ("Delivers the brief's explicit asks:", 0, True),
    "Central monitoring (Anypoint Monitoring/Visualizer)",
    "Reusable connection templates (Exchange + API-led)",
    "Admin error queue (DLQ / reprocessing)",
    "Uniform mutual TLS + security policy (API Manager)",
    ("Pays for itself:", 0, True),
    "Reuse \u2192 faster go-to-market; new agency/system = config, not code",
    "Lower TCO; fewer governor-limit issues; audit-ready (regulated FS)",
    "Salesforce-native product (Anypoint + Composer for clicks-not-code)",
], size=12.5, gap=6)

box(s, Inches(6.9), Inches(1.25), Inches(6.0), Inches(0.45), "Key trade-offs articulated", NAVY, font=14, align=PP_ALIGN.LEFT)
bullets(s, Inches(6.95), Inches(1.8), Inches(5.9), Inches(3.0), [
    "Virtualize vs replicate history \u2192 Salesforce Connect + cache hot path (don't store 100M+/day)",
    "Sync vs async credit checks \u2192 Apex Continuation to scale 30 concurrent \u00d7 1 min",
    "Platform Events vs CDC \u2192 CDC for INT-19 (native history + changed-fields + deletes)",
    "Native vs middleware \u2192 Mule where monitoring/reuse/error-queue pays; native for pure UI & atomic writes",
], size=12.5, gap=8)
box(s, Inches(6.9), Inches(4.95), Inches(6.0), Inches(0.45), "Quantified outcomes", NAVY, font=14, align=PP_ALIGN.LEFT)
bullets(s, Inches(6.95), Inches(5.5), Inches(5.9), Inches(1.7), [
    "Scales: 5M customers / ~100M txns/day / >100k history calls/hr / millions of loan events/day",
    "Resilient (replay + DLQ) \u00b7 Secure (mTLS + end-user auth) \u00b7 Reusable (API-led assets)",
], size=12.5, gap=8)

# ============================================================
# THANK YOU
# ============================================================
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
textbox(s, Inches(0), Inches(3.0), SW, Inches(1.5),
        [("Thank you", 48, True, WHITE), ("Questions?", 20, False, RGBColor(0x7F, 0xB8, 0xEE))],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

out = "Bedrock - Integration Architecture - Presentation.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")

